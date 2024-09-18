import os
import json
import openai
from flask import Flask, render_template, request, jsonify
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader
import mammoth
from pptx import Presentation
from docx import Document

# Initialize Flask app
app = Flask(__name__)

# Configure upload folder and allowed extensions
UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Ensure upload folder exists
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# Allowed file extensions
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'pptx', 'txt'}

# OpenAI API key (replace with your actual key)
openai.api_key = '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'

# Check if file is allowed
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# Read PDF file
def read_pdf(file_path):
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

# Read DOCX file
def read_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

# Read PPTX file
def read_pptx(file_path):
    presentation = python_pptx.Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Read TXT file
def read_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

# Home route
@app.route('/')
def home():
    return render_template('index.html')

# File upload route
@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'success': False, 'message': 'No file part'})
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'success': False, 'message': 'No selected file'})
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        # Read the content based on the file type
        if filename.endswith('.pdf'):
            file_content = read_pdf(file_path)
        elif filename.endswith('.docx'):
            file_content = read_docx(file_path)
        elif filename.endswith('.pptx'):
            file_content = read_pptx(file_path)
        elif filename.endswith('.txt'):
            file_content = read_txt(file_path)
        else:
            return jsonify({'success': False, 'message': 'Unsupported file type'})
        
        return jsonify({'success': True, 'content': file_content})
    
    return jsonify({'success': False, 'message': 'File type not allowed'})

# Route for handling AI questions
@app.route('/ask', methods=['POST'])
def ask_question():
    data = request.get_json()
    file_content = data.get('content')
    question = data.get('question')
    
    if not file_content or not question:
        return jsonify({'success': False, 'message': 'Invalid input'})

    # Send the file content and question to OpenAI
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt=f"The document contains the following text:\n{file_content}\n\nQ: {question}\nA:",
        temperature=0.5,
        max_tokens=150
    )

    answer = response.choices[0].text.strip()
    return jsonify({'success': True, 'answer': answer})

if __name__ == '__main__':
    app.run(debug=True)
